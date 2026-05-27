"""PPTX builder: constructs slides with glass blur background,
fold-cropped image, and decorative gradient triangles.

Uses direct lxml manipulation on python-pptx elements (no zipfile post-processing).
Dynamically adjusts fold geometry and triangle positions based on image aspect ratio."""

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

# Slide dimensions (16:9 widescreen)
SLIDE_W = 12192000
SLIDE_H = 6858000

# Namespace URIs
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _qname(ns, tag):
    return "{%s}%s" % (ns, tag)


def build_pptx(image_specs, output_path):
    """Build a PPTX with one slide per image.

    Args:
        image_specs: list of dicts with keys:
            - orig_path: path to original image
            - glass_path: path to glass-blurred image (Pillow pre-computed)
            - img_width: original image width in pixels
            - img_height: original image height in pixels
        output_path: path for the output .pptx file
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for spec in image_specs:
        slide = prs.slides.add_slide(blank_layout)

        # Remove placeholder shapes
        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)

        _build_slide(slide, spec)

    prs.save(output_path)


def _build_slide(slide, spec):
    """Build one slide: glass background + fold image + triangle group."""
    orig_path = spec["orig_path"]
    glass_path = spec["glass_path"]
    iw = spec["img_width"]
    ih = spec["img_height"]

    # Calculate image dimensions when scaled to slide height (aspect ratio preserved)
    img_cx = int(SLIDE_H * iw / ih)   # image width at full slide height
    img_x = SLIDE_W - img_cx           # right-aligned

    # Delta = difference in image width from reference (square image)
    REFERENCE_IMG_CX = 6858000
    delta = img_cx - REFERENCE_IMG_CX

    # 1. Add glass-blurred background (full slide, plain picture, no effects)
    bg_pic = slide.shapes.add_picture(glass_path, 0, 0, SLIDE_W, SLIDE_H)
    bg_pic.name = "_bg_glass_"

    # 2. Add original image (full height, aspect ratio preserved, right-aligned)
    fold_pic = slide.shapes.add_picture(
        orig_path, Emu(img_x), Emu(0), Emu(img_cx), Emu(SLIDE_H)
    )
    fold_pic.name = "_fold_image_"

    # 3. Inject custom fold geometry — returns fold crease position on slide
    crease_slide_x = _inject_fold_geometry(fold_pic._element, img_cx, delta)

    # 4. Add decorative triangle group at the fold crease
    _add_triangle_group(slide._element, crease_slide_x)


def _inject_fold_geometry(pic_element, img_cx, delta):
    """Replace the fold picture's preset geometry with the '多边形' mask polygon.

    Uses the delta-based approach:
    - Rightmost points (at the slide right edge) keep their coordinates
    - Fold crease points shift by -delta (delta = img_cx - reference_img_cx)
    - Then normalize so all x >= 0 (shift globally if needed)
    """
    spPr = pic_element.find(".//{%s}spPr" % NS_P)
    if spPr is None:
        return

    # Remove existing preset geometry
    for prstGeom in spPr.findall(_qname(NS_A, "prstGeom")):
        spPr.remove(prstGeom)

    # --- Original polygon points ('多边形' from slide2) ---
    # Path dimensions: w=6421120, h=6878320 (extends 20320 EMU above slide)
    # Rightmost x values (anchored to slide right edge, DON'T shift by delta):
    RIGHTMOST_X = {6421120, 6417733, 6414347, 6410960}

    orig_points = [
        ("moveTo",     0, 3454400),
        ("lnTo",       3484880, 6858000),
        ("lnTo",       6421120, 6878320),
        ("cubicBezTo", 6417733, 4592320, 6414347, 2306320, 6410960, 20320),
        ("lnTo",       3464560, 0),
        ("lnTo",       71120, 3393440),
        ("lnTo",       0, 3454400),
        ("close",),
    ]

    # --- Shift movable x by -delta, cap y to slide height ---
    Y_CAP = SLIDE_H  # 6858000
    shifted = []
    for pt in orig_points:
        cmd = pt[0]
        if cmd == "close":
            shifted.append(("close",))
        elif cmd == "cubicBezTo":
            def _shift_x(val):
                return val if val in RIGHTMOST_X else int(val - delta)
            shifted.append(("cubicBezTo",
                _shift_x(pt[1]), min(pt[2], Y_CAP),
                _shift_x(pt[3]), min(pt[4], Y_CAP),
                _shift_x(pt[5]), min(pt[6], Y_CAP)))
        else:
            x = pt[1] if pt[1] in RIGHTMOST_X else int(pt[1] - delta)
            y = min(pt[2], Y_CAP)
            shifted.append((cmd, x, y))

    # --- Normalize: ensure all x >= 0 ---
    all_x = []
    for pt in shifted:
        if pt[0] in ("moveTo", "lnTo"):
            all_x.append(pt[1])
        elif pt[0] == "cubicBezTo":
            all_x.extend([pt[1], pt[3], pt[5]])

    min_x = min(all_x)
    if min_x < 0:
        shift_amount = -min_x
        normalized = []
        for pt in shifted:
            if pt[0] == "close":
                normalized.append(("close",))
            elif pt[0] == "cubicBezTo":
                normalized.append(("cubicBezTo",
                    pt[1] + shift_amount, pt[2],
                    pt[3] + shift_amount, pt[4],
                    pt[5] + shift_amount, pt[6]))
            else:
                normalized.append((pt[0], pt[1] + shift_amount, pt[2]))
        shifted = normalized

    # Recalculate path dimensions after normalization
    all_x = []
    for pt in shifted:
        if pt[0] in ("moveTo", "lnTo"):
            all_x.append(pt[1])
        elif pt[0] == "cubicBezTo":
            all_x.extend([pt[1], pt[3], pt[5]])

    path_w = max(all_x)
    path_h = Y_CAP  # capped to slide height

    # Update fold picture dimensions to match path
    fold_x = SLIDE_W - path_w
    # Also adjust y to fit path: polygon extends above slide by 20320 EMU,
    # so position picture at y=-20320 with cy=path_h to cover full path area
    fold_y = 0
    fold_cy = path_h
    xfrm = spPr.find(_qname(NS_A, "xfrm"))
    if xfrm is not None:
        off = xfrm.find(_qname(NS_A, "off"))
        if off is not None:
            off.set("x", str(fold_x))
            off.set("y", str(fold_y))
        ext = xfrm.find(_qname(NS_A, "ext"))
        if ext is not None:
            ext.set("cx", str(path_w))
            ext.set("cy", str(fold_cy))

    # --- Build custGeom ---
    cg = etree.SubElement(spPr, _qname(NS_A, "custGeom"))
    etree.SubElement(cg, _qname(NS_A, "avLst"))

    # gdLst — formula values match the shifted+normalized coordinates
    gdLst_data = [
        ("csX0", "*/ 0 w %d" % path_w),       ("csY0", "*/ 3454400 h %d" % path_h),
        ("csX1", "*/ 3484880 w %d" % path_w), ("csY1", "*/ 6858000 h %d" % path_h),
        ("csX2", "*/ 6421120 w %d" % path_w), ("csY2", "*/ 6878320 h %d" % path_h),
        ("csX3", "*/ 6417733 w %d" % path_w), ("csY3", "*/ 4592320 h %d" % path_h),
        ("csX4", "*/ 6414347 w %d" % path_w), ("csY4", "*/ 2306320 h %d" % path_h),
        ("csX5", "*/ 6410960 w %d" % path_w), ("csY5", "*/ 20320 h %d" % path_h),
        ("csX6", "*/ 3464560 w %d" % path_w), ("csY6", "*/ 0 h %d" % path_h),
        ("csX7", "*/ 71120 w %d" % path_w),   ("csY7", "*/ 3393440 h %d" % path_h),
        ("csX8", "*/ 0 w %d" % path_w),       ("csY8", "*/ 3454400 h %d" % path_h),
    ]
    gdLst = etree.SubElement(cg, _qname(NS_A, "gdLst"))
    for name, fmla in gdLst_data:
        gd = etree.SubElement(gdLst, _qname(NS_A, "gd"))
        gd.set("name", name)
        gd.set("fmla", fmla)

    etree.SubElement(cg, _qname(NS_A, "ahLst"))

    # cxnLst
    cxnLst = etree.SubElement(cg, _qname(NS_A, "cxnLst"))
    for i in range(9):
        cxn = etree.SubElement(cxnLst, _qname(NS_A, "cxn"))
        cxn.set("ang", "0")
        pos = etree.SubElement(cxn, _qname(NS_A, "pos"))
        pos.set("x", f"csX{i}")
        pos.set("y", f"csY{i}")

    # rect
    rect = etree.SubElement(cg, _qname(NS_A, "rect"))
    rect.set("l", "l"); rect.set("t", "t"); rect.set("r", "r"); rect.set("b", "b")

    # pathLst — direct coordinates
    pathLst = etree.SubElement(cg, _qname(NS_A, "pathLst"))
    path = etree.SubElement(pathLst, _qname(NS_A, "path"))
    path.set("w", str(path_w))
    path.set("h", str(path_h))

    for pt in shifted:
        cmd = pt[0]
        if cmd == "moveTo":
            el = etree.SubElement(path, _qname(NS_A, "moveTo"))
            sub = etree.SubElement(el, _qname(NS_A, "pt"))
            sub.set("x", str(pt[1])); sub.set("y", str(pt[2]))
        elif cmd == "lnTo":
            el = etree.SubElement(path, _qname(NS_A, "lnTo"))
            sub = etree.SubElement(el, _qname(NS_A, "pt"))
            sub.set("x", str(pt[1])); sub.set("y", str(pt[2]))
        elif cmd == "cubicBezTo":
            el = etree.SubElement(path, _qname(NS_A, "cubicBezTo"))
            for j in range(3):
                sub = etree.SubElement(el, _qname(NS_A, "pt"))
                sub.set("x", str(pt[1 + j*2])); sub.set("y", str(pt[2 + j*2]))
        elif cmd == "close":
            etree.SubElement(path, _qname(NS_A, "close"))

    # Fold crease slide position: the leftmost movable polygon point
    # shifts by -delta (where delta = img_cx - 6858000).
    # Any normalization shifts fold_x by the same amount, so the crease
    # slide position simplifies to: original_fold_x - delta = 5770880 - delta.
    crease_slide_x = 5770880 - delta
    return crease_slide_x


def _add_triangle_group(slide_element, crease_slide_x):
    """Add the decorative gradient triangle group to the slide's shape tree.
    Position is based on the fold crease slide position."""
    nsmap = {"p": NS_P, "a": NS_A}
    spTree = slide_element.find(".//{%s}cSld/{%s}spTree" % (NS_P, NS_P))
    if spTree is None:
        return

    # Sample: fold crease at slide_x = 5770880, group at slide_x = 6096000
    # Offset = 6096000 - 5770880 = 325120 EMU to the right of the crease
    group_x = crease_slide_x + 325120
    group_y = -755030  # unchanged

    grpSp = etree.SubElement(spTree, _qname(NS_P, "grpSp"))

    # Group non-visual properties
    nvGrpSpPr = etree.SubElement(grpSp, _qname(NS_P, "nvGrpSpPr"))
    cNvPr = etree.SubElement(nvGrpSpPr, _qname(NS_P, "cNvPr"))
    cNvPr.set("id", "16")
    cNvPr.set("name", "Group 15")
    etree.SubElement(nvGrpSpPr, _qname(NS_P, "cNvGrpSpPr"))
    etree.SubElement(nvGrpSpPr, _qname(NS_P, "nvPr"))

    # Group transform
    grpSpPr = etree.SubElement(grpSp, _qname(NS_P, "grpSpPr"))
    xfrm = etree.SubElement(grpSpPr, _qname(NS_A, "xfrm"))
    off = etree.SubElement(xfrm, _qname(NS_A, "off"))
    off.set("x", str(group_x))
    off.set("y", str(group_y))
    ext_off = etree.SubElement(xfrm, _qname(NS_A, "ext"))
    ext_off.set("cx", "1918089")
    ext_off.set("cy", "8368059")
    chOff = etree.SubElement(xfrm, _qname(NS_A, "chOff"))
    chOff.set("x", "5914194")
    chOff.set("y", "-1838587")
    chExt = etree.SubElement(xfrm, _qname(NS_A, "chExt"))
    chExt.set("cx", "1918089")
    chExt.set("cy", "9684705")

    # Bottom triangle
    _add_triangle_shape(grpSp, "14", "Isosceles Triangle 13",
                        x=5914195, y=2587465, cx=1918088, cy=5258653,
                        rot=19576042, flipV=0)

    # Top triangle (flipped vertically)
    _add_triangle_shape(grpSp, "15", "Isosceles Triangle 13",
                        x=5914194, y=-1838587, cx=1918088, cy=5258653,
                        rot=2023958, flipV=1)


def _add_triangle_shape(parent, cNvPr_id, name, x, y, cx, cy, rot, flipV):
    """Add a single triangle shape with gradient fill to the group."""
    sp = etree.SubElement(parent, _qname(NS_P, "sp"))

    # Non-visual properties
    nvSpPr = etree.SubElement(sp, _qname(NS_P, "nvSpPr"))
    cNvPr_el = etree.SubElement(nvSpPr, _qname(NS_P, "cNvPr"))
    cNvPr_el.set("id", cNvPr_id)
    cNvPr_el.set("name", name)
    etree.SubElement(nvSpPr, _qname(NS_P, "cNvSpPr"))
    etree.SubElement(nvSpPr, _qname(NS_P, "nvPr"))

    # Shape properties
    spPr = etree.SubElement(sp, _qname(NS_P, "spPr"))

    # Transform
    xfrm = etree.SubElement(spPr, _qname(NS_A, "xfrm"))
    xfrm.set("rot", str(rot))
    if flipV:
        xfrm.set("flipV", "1")
    off_el = etree.SubElement(xfrm, _qname(NS_A, "off"))
    off_el.set("x", str(x))
    off_el.set("y", str(y))
    ext_el = etree.SubElement(xfrm, _qname(NS_A, "ext"))
    ext_el.set("cx", str(cx))
    ext_el.set("cy", str(cy))

    # Custom geometry
    custGeom = etree.SubElement(spPr, _qname(NS_A, "custGeom"))

    # avLst + gdLst
    etree.SubElement(custGeom, _qname(NS_A, "avLst"))

    tri_gd_data = [
        ("csX0", "*/ 0 w 1918088"),       ("csY0", "*/ 3962400 h 5258653"),
        ("csX1", "*/ 939800 w 1918088"),  ("csY1", "*/ 0 h 5258653"),
        ("csX2", "*/ 1918088 w 1918088"), ("csY2", "*/ 5258653 h 5258653"),
        ("csX3", "*/ 0 w 1918088"),       ("csY3", "*/ 3962400 h 5258653"),
    ]
    gdLst = etree.SubElement(custGeom, _qname(NS_A, "gdLst"))
    for gd_name, gd_fmla in tri_gd_data:
        gd = etree.SubElement(gdLst, _qname(NS_A, "gd"))
        gd.set("name", gd_name)
        gd.set("fmla", gd_fmla)

    etree.SubElement(custGeom, _qname(NS_A, "ahLst"))

    cxnLst = etree.SubElement(custGeom, _qname(NS_A, "cxnLst"))
    for i in range(4):
        cxn = etree.SubElement(cxnLst, _qname(NS_A, "cxn"))
        cxn.set("ang", "0")
        pos = etree.SubElement(cxn, _qname(NS_A, "pos"))
        pos.set("x", f"csX{i}")
        pos.set("y", f"csY{i}")

    rect = etree.SubElement(custGeom, _qname(NS_A, "rect"))
    rect.set("l", "l"); rect.set("t", "t"); rect.set("r", "r"); rect.set("b", "b")

    pathLst = etree.SubElement(custGeom, _qname(NS_A, "pathLst"))
    path = etree.SubElement(pathLst, _qname(NS_A, "path"))
    path.set("w", "1918088")
    path.set("h", "5258653")

    tri_points = [
        ("moveTo", 0, 3962400),
        ("lnTo", 939800, 0),
        ("lnTo", 1918088, 5258653),
        ("cubicBezTo", 1624625, 4503708, 1618143, 4343552, 0, 3962400),
        ("close",),
    ]
    for pt in tri_points:
        cmd = pt[0]
        if cmd == "moveTo":
            el = etree.SubElement(path, _qname(NS_A, "moveTo"))
            sub = etree.SubElement(el, _qname(NS_A, "pt"))
            sub.set("x", str(pt[1])); sub.set("y", str(pt[2]))
        elif cmd == "lnTo":
            el = etree.SubElement(path, _qname(NS_A, "lnTo"))
            sub = etree.SubElement(el, _qname(NS_A, "pt"))
            sub.set("x", str(pt[1])); sub.set("y", str(pt[2]))
        elif cmd == "cubicBezTo":
            el = etree.SubElement(path, _qname(NS_A, "cubicBezTo"))
            for j in range(3):
                sub = etree.SubElement(el, _qname(NS_A, "pt"))
                sub.set("x", str(pt[1 + j*2])); sub.set("y", str(pt[2 + j*2]))
        elif cmd == "close":
            etree.SubElement(path, _qname(NS_A, "close"))

    # Gradient fill
    gradFill = etree.SubElement(spPr, _qname(NS_A, "gradFill"))
    gsLst = etree.SubElement(gradFill, _qname(NS_A, "gsLst"))

    gs0 = etree.SubElement(gsLst, _qname(NS_A, "gs"))
    gs0.set("pos", "0")
    schClr0 = etree.SubElement(gs0, _qname(NS_A, "schemeClr"))
    schClr0.set("val", "accent1")
    etree.SubElement(schClr0, _qname(NS_A, "lumMod")).set("val", "5000")
    etree.SubElement(schClr0, _qname(NS_A, "lumOff")).set("val", "95000")

    gs50 = etree.SubElement(gsLst, _qname(NS_A, "gs"))
    gs50.set("pos", "50000")
    schClr50 = etree.SubElement(gs50, _qname(NS_A, "schemeClr"))
    schClr50.set("val", "bg1")
    etree.SubElement(schClr50, _qname(NS_A, "lumMod")).set("val", "85000")

    gs100 = etree.SubElement(gsLst, _qname(NS_A, "gs"))
    gs100.set("pos", "100000")
    schClr100 = etree.SubElement(gs100, _qname(NS_A, "schemeClr"))
    schClr100.set("val", "bg1")

    lin = etree.SubElement(gradFill, _qname(NS_A, "lin"))
    lin.set("ang", "10800000")
    lin.set("scaled", "0")

    ln = etree.SubElement(spPr, _qname(NS_A, "ln"))
    etree.SubElement(ln, _qname(NS_A, "noFill"))

    # Empty text body (center aligned)
    txBody = etree.SubElement(sp, _qname(NS_P, "txBody"))
    bodyPr = etree.SubElement(txBody, _qname(NS_A, "bodyPr"))
    bodyPr.set("rtlCol", "0")
    bodyPr.set("anchor", "ctr")
    etree.SubElement(txBody, _qname(NS_A, "lstStyle"))
    ap = etree.SubElement(txBody, _qname(NS_A, "p"))
    apPr = etree.SubElement(ap, _qname(NS_A, "pPr"))
    apPr.set("algn", "ctr")
    endRPr = etree.SubElement(ap, _qname(NS_A, "endParaRPr"))
    endRPr.set("lang", "zh-CN")
    endRPr.set("altLang", "en-US")
